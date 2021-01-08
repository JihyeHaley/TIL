import mammoth
import glob
import re
import os
import xlsxwriter
import timeit
from itertools import zip_longest

# help function
# 한글 / 영어 둘다 아닌 경우
def neitherKoNorEn(text):
    neither = re.compile('((?![ㄱ-ㅣ가-힣|a-zA-Z]).)*')
    return bool(neither.fullmatch(text))


# 한글이 포함되지 않은 경우
def mustEnglish(text):
    never_ko = re.compile('((?![ㄱ-ㅣ가-힣]).)*')
    return bool(never_ko.fullmatch(text))


# 영어가 포함되어있는 경우
def isEnglish(text):
    en = re.compile('.*[a-zA-Z]+')
    # return bool(ko.fullmatch(text))
    return bool(en.match(text))


# 한글이 포함되어있는 경우
def isKorean(text):
    ko = re.compile('.*[ㄱ-ㅣ가-힣]+')
    # return bool(ko.fullmatch(text))
    return bool(ko.match(text))



# mustEnglish - false & isKorean - true
# => 한글만 이거나 한글영어 섞여있을때
# 이경우 다시 sEnglish로 다시 check
# false면 한글만 혹은.. 영어없이 숫자나 특수문자만
# true면 영어+한글 혼합

# mustEnglish - true & isKorean - false
# => 영어만

# --------------------------------------------------------------------------------------
# def sentence_tokenizer(text, source_language_code):
#     tokenizer = stanza.Pipeline(lang=source_language_code, processors='tokenize', verbose=False)
#     doc = tokenizer(text)
#     return [sentence.text for sentence in doc.sentences]


def is_skip(text):
    skip = re.compile('^[0-9『』@%#^&~!★●▪:;()/=+*$,._\\\\-\\s]+')
    return bool(skip.fullmatch(text))


def is_page_num(text):
    is_page = re.compile('.[0-9|pP[\\]{}()<>\\s]+')
    return bool(is_page.fullmatch(text))


def get_excel_index(lang, row_idx):
    if lang == 'ko':
        return f'A{str(row_idx)}'
    else:
        return f'B{str(row_idx)}'


# ppt인 파일 가져오기

from pptx import Presentation           
    
# dictionary로 넣어주기
file_lists_ko = {}
file_lists_en = {}
file_lists_mer = {}   


# dictionary로 넣어주기
file_lists_ko = {}
file_lists_en = {}
file_lists_mer = {}   

def get_filename_list(path, ext):
    return [f for f in glob.glob(path + f"/*{ext}")]

root_path = '/Users/haley/Desktop/Farm'
sub_dir = '/4_2019한국표준협회/4.2019한국표준협회'
file_path = root_path + sub_dir

def file_to_excel():
    start = timeit.default_timer()
    # -----------------------------------------------------------------------
    # file 다 져오기

    # file_name_lists = [f for f in glob.glob(f'./4.2019한국표준협회/*.*')]
    docx_name_lists = get_filename_list(file_path, '.docx')
    pptx_name_lists = get_filename_list(file_path, '.pptx')

    doc_name_lists = get_filename_list(file_path, '.doc')
    rtf_name_lists = get_filename_list(file_path, '.rtf')
    txt_name_lists = get_filename_list(file_path, '.txt')

    # etc_name_lists = []
    
    print(file_path + f' docx: {len(docx_name_lists)}, pptx: {len(pptx_name_lists)}, doc: {len(doc_name_lists)}')
    
    
    # ------------------------------------------------------------------------------------
    # ppt dict 만들기
    from pptx import Presentation

    file_lists_pptx = {}
    i = 0

    # dict에 잘 들어갔는지 test
    # pptx_test_print_1 = ''
    # pptx_test_print_2 = ''
    # pptx_test_print_s = ''


    # ppt 추출 시작
    for pptx_name_list in pptx_name_lists:
        i += 1
        # 일단 다 긁어오기 위한 list
        pptx_results_pre = []

        # 최종 잘 나눠진 아이를 넣을 1개 pptx의 list
        pptx_results = []
        
        # pptx 분석 하기위해 넣어주기
        prs = Presentation(pptx_name_list)
        
        # 1개의 pptx 분석시작
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text == '':
                        continue
                    pptx_results_pre.append(paragraph.text)


        for pptx_result_pre in pptx_results_pre:
            if type(pptx_result_pre) == 'list':
                for line in pptx_result_pre:
                    pptx_results.append(line.split())
            else:
                pptx_results.append(pptx_result_pre)
        pptx_results_pre = list(filter(None, pptx_results_pre)) 

        # ----------------------------------------------------
        # print test
        # pptx_name_list의 길이가 너무 깁니다.
        # '/Users/haley/Desktop/Farm/4_2019한국표준협회/4.2019한국표준협회/TE-한국표준협회-ISOfocus_131_en-한.pptx'
        # 그래서 [38:-5]로 slicing 했습니다.
        # -5는 뒤에 .docx, .pptx 를 자른겁니다!
        if i == 1:
            pptx_test_print_1 = pptx_name_list[38:-5]
            pptx_test_print_s = pptx_name_list
        elif i == 2:
            pptx_test_print_2 = pptx_name_list[38:-5]
        else:
            print(i, end=' ')

        # dict에 들어갈 key값을 줄여줬습니다.
        file_lists_pptx[pptx_name_list[38:-5]] = pptx_results


    print(f'done, file_lists_pptx : is {len(file_lists_pptx)}')
    # print(file_lists_pptx.get(pptx_test_print_1))
    # print(file_lists_pptx.get(pptx_test_print_2))
    # print(f'done, file_full name : is {pptx_test_print_s}')
    
    
    
    # ------------------------------------------------------------------------------------
    # docx dict 만들기
    file_lists_docx = {}
    i = 0
    
    # dict에 잘 들어갔는지 test
    # docx_test_print_1 = ''
    # docx_test_print_2 = ''
    
    # docx 추출 시작
    for docx_name_list in docx_name_lists:
        i += 1
        docx_results = []

        with open(docx_name_list, "rb") as docx_file:
            result = mammoth.extract_raw_text(docx_file)
            text = result.value  # The raw text
            contents = text.split('\n |" "')
    
        for line in contents:
            line.replace(".", ". ")
        docx_results = text.split('\n')
        docx_results = list(filter(None, docx_results))


        # ----------------------------------------------------
        # print test
        # pptx_name_list의 길이가 너무 깁니다.
        # '/Users/haley/Desktop/Farm/4_2019한국표준협회/4.2019한국표준협회/TE-한국표준협회-ISOfocus_131_en-한.pptx'
        # 그래서 [38:-5]로 slicing 했습니다.
        # -5는 뒤에 .docx, .pptx 를 자른겁니다!
        if i == 1:
            docx_test_print_1 = docx_name_list[38:-5]
        elif i == 2:
            docx_test_print_2 = docx_name_list[38:-5]
        else:
            print(i, end=' ')

        # dict에 들어갈 key값을 줄여줬습니다.
        file_lists_docx[docx_name_list[38:-5]] = docx_results

    print(f'done, file_lists_docx : is {len(file_lists_docx)}')
    # print(file_lists_docx.get(docx_test_print_1))
    # print(file_lists_docx.get(docx_test_print_2))
    


    # ------------------------------------------------------------------------------------
    # rtf dict 만들기
    file_lists_rtf = {}
    i = 0
    
    # dict에 잘 들어갔는지 test
    rtf_test_print_1 = ''
    rtf_test_print_sample = ''
    
    # rtf 추출 시작
    for rtf_name_list in rtf_name_lists:
        i += 1
        rtf_results = []
    
        with open(rtf_name_list, 'r') as rtf_file:
            rtf_results = rtf_file.readline()
            
        
        # rtf_results = list(filter(None, rtf_results))

        # ----------------------------------------------------
        # print test
        # pptx_name_list의 길이가 너무 깁니다.
        # '/Users/haley/Desktop/Farm/4_2019한국표준협회/4.2019한국표준협회/TE-한국표준협회-ISOfocus_131_en-한.pptx'
        # 그래서 [38:-5]로 slicing 했습니다.
        # -5는 뒤에 .docx, .pptx 를 자른겁니다!
        if i == 1:
            rtf_test_print_1 = rtf_name_list[38:-4]
            rtf_test_print_sample = rtf_name_list[38:-4]
        else:
            print(i, end=' ')

        # dict에 들어갈 key값을 줄여줬습니다.
        file_lists_rtf[rtf_name_list[38:-4]] = rtf_results

    print(f'done, file_lists_txt : is {len(file_lists_rtf)}')
    print(file_lists_rtf.get(rtf_test_print_1))
    print(rtf_test_print_sample)



    # ------------------------------------------------------------------------------------
    # txt dict 만들기
    file_lists_txt = {}
    i = 0
    
    # dict에 잘 들어갔는지 test
    txt_test_print_1 = ''
    txt_test_print_sample = ''
    
    # txt 추출 시작
    for txt_name_list in txt_name_lists:
        i += 1
        txt_results = []
    
        with open(txt_name_list, 'r') as txt_file:
            txt_results = txt_file.readline()
            
        # for txt_result in txt_results:
        #     txt_results.append(txt_result.strip())
        # txt_results = list(filter(None, txt_results))

        # ----------------------------------------------------
        # print test
        # pptx_name_list의 길이가 너무 깁니다.
        # '/Users/haley/Desktop/Farm/4_2019한국표준협회/4.2019한국표준협회/TE-한국표준협회-ISOfocus_131_en-한.pptx'
        # 그래서 [38:-5]로 slicing 했습니다.
        # -5는 뒤에 .docx, .pptx 를 자른겁니다!
        if i == 1:
            txt_test_print_1 = txt_name_list[38:-4]
            txt_test_print_sample = txt_name_list[38:-4]
        else:
            print(i, end=' ')

        # dict에 들어갈 key값을 줄여줬습니다.
        file_lists_txt[txt_name_list[38:-4]] = txt_results

    print(f'done, file_lists_txt : is {len(file_lists_txt)}')
    print(file_lists_txt.get(txt_test_print_1))
    print(txt_test_print_sample)
    
    
    
    # ------------------------------------------------------------------------------------
    # docx와, pptx merging
    file_lists_all = {}
    file_lists_all.update(file_lists_pptx)
    file_lists_all.update(file_lists_docx)
    file_lists_all.update(file_lists_txt)
    file_lists_all.update(file_lists_rtf)
    print(f'total length is : {len(file_lists_all)}')
    # print(list(file_lists_all.keys()))
    
    
    # ------------------------------------------------------------------------------------
    # 한글, 영어 구분해서 각각의 dict에 넣어주기
    ko_files = {}
    en_files = {}
    
    i = 0
    j = 1
    for key, value in file_lists_all.items():
        i += 1 
        # 언어 type을 구분하기 위한 slicing 입니다.
        language = key[-3:]

        # key_name = '/4.2019한국표준협회/TE-한국표준협회-ISOfocus_131_en-한.pptx'
        # '/4.' 를 제외하고 확장자명(.pptx, .docx)를 제외하기 위한 slicing입니다.
        key_name = key[3:-4]
        print(key[-3:]) # check!
        
        if language == '한':
            ko_files.update({key_name : value})
            print(key_name) # 잘 들어갔는지 체크하기위해서
        elif language == '영':
            en_files.update({key_name : value})
            print(key_name) # 잘 들어갔는지 체크하기위해서
        elif language == '병':
            j += 1 
    
    print(j)
    print(list(ko_files.keys()))
    print(list(en_files.keys()))
     
    

    
    # ------------------------------------------------------------------------------------
    # excel 밖으로 빼내기
    workbook = xlsxwriter.Workbook('/Users/haley/Desktop/Farm/4_2019한국표준협회/standard6_.xlsx')
    print('excel 빼내기 파일생성 성공!')
    worksheet = workbook.add_worksheet()
    worksheet.write('A1', 'ko')
    worksheet.write('B1', 'en')
    
    row_idx = 2
    
    i = 0
    

    # ------------------------------------------------------------------------------------
    # excel 글쓰기
    for ko_key, ko_value in ko_files.items():
        index = get_excel_index('ko', row_idx)
        # worksheet.write(f'A{str(row_idx)}', '>'*20 + ko_key)
        worksheet.write(index, '>'*20 + ko_key)
        for en_key, en_value in en_files.items():
            if ko_key == en_key:
                index = get_excel_index('en', row_idx)
                worksheet.write(index, '>'*20 + en_key)
                for ko, en in zip_longest(ko_value, en_value, fillvalue=' '):
                    ko_index = get_excel_index('ko', row_idx)
                    en_index = get_excel_index('en', row_idx)
                    worksheet.write(ko_index, ko)
                    worksheet.write(en_index, en)
                    row_idx += 1
            
    workbook.close()
    stop = timeit.default_timer()
    
    print(' Running Time: ', stop - start)

file_to_excel()
i
