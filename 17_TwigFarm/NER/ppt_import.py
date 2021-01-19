# file 다 가져오기
import glob
import subprocess
from pptx import Presentation    

lower = [chr(l) for l in range(97, 123, 1)]
upper = [chr(u) for u in range(65, 91, 1)] 

def get_filename_list(path, ext):
    print('Change encoding files.......\ \n')
    result = []
    for f in glob.glob(file_path + f"/*{ext}"):
        # run convmv shell script -> file normalization NFC -> NFD (한글자소분리해결)
        subprocess.run(['/usr/local/bin/convmv', '-f', 'utf-8', '-t', 'utf-8', '--nfc', '--notest', f])
        result.append(f)
    return result

root_path = '/Users/jihyeoh/Lexcode/팀 채널 - 인공지능 학습 DB 구축 채널/1. 원본DB(렉스코드)'
sub_dir = '/4.2019한국표준협회'

file_path = root_path + sub_dir
pptx_name_lists = get_filename_list(file_path, 'TE-한국표준협회-ISOfocus_131_en-영.pptx')

for i in pptx_name_lists:
    prs = Presentation(i)
    result = []
    for slide in prs.slides:
        for shape in slide.shapes:

            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                print(f'{paragraph.text}')
                if paragraph.text == '':
                    continue
                # if  == True:
                #     continue
                result.append(paragraph.text)
    



pptx_results = []

for idx in range(len(result)):
    stop_lower_idx = 0 
    found_upper_idx = 0
    finally_sentenced = ''
    print(result[idx])
    if result[idx][:3] == 'ver':
        stop_lower_idx = idx
        for jdx in range(idx, len(result)):
            if result[jdx][0] == 'E':
                found_upper_idx = jdx
                finally_sentenced = result[found_upper_idx] + result[stop_lower_idx]

    if idx == stop_lower_idx:
        pptx_results.append(finally_sentenced)
    elif idx == found_upper_idx:
        continue
    elif len(result[idx]) == 1 and result[idx][0] in upper:
        continue
    else:
        pptx_results.append(result[idx])
    
pptx_results = list(filter(None, pptx_results))

for pptx_result in pptx_results:
    print(pptx_result, end='\n\n')