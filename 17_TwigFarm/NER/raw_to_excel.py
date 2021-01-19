# -----------------------------------------------------------------------------------------
# library import!!!
import mammoth
import glob
import re
import os
import xlsxwriter
import timeit
import sys
from itertools import zip_longest
from pptx import Presentation
from nltk.tokenize import sent_tokenize
import subprocess
import nltk
# ppt인 파일 가져오기
from pptx import Presentation    


# -----------------------------------------------------------------------------------------
# pattern의 regex를 text에서 찾아 지움
def remove_regex(pattern, text):
    if re.search(pattern, text):
        return re.sub(pattern, '', text)
    return text


# -----------------------------------------------------------------------------------------
# excel numbering
def get_excel_index(lang, row_idx):
    if lang == 'ko':
        return f'B{str(row_idx)}'
    else:
        return f'C{str(row_idx)}'


# -----------------------------------------------------------------------------------------
# file 다 가져오기
def get_filename_list(path, ext):
    print('Change encoding files.......\ \n')
    result = []
    for f in glob.glob(file_path + f"/*{ext}"):
        # run convmv shell script -> file normalization NFC -> NFD (한글자소분리해결)
        subprocess.run(['/usr/local/bin/convmv', '-f', 'utf-8', '-t', 'utf-8', '--nfc', '--notest', f])
        result.append(f)
    return result

root_path = '/Users/jihyeoh/Lexcode/팀 채널 - 인공지능 학습 DB 구축 채널/1. 원본DB(렉스코드)'
sub_dir = '/4.2019한국표준협회'
file_path = root_path + sub_dir


# -----------------------------------------------------------------------------------------
# 대소문자 구별하기 위한 
def lower_or_upper():
    lower = [chr(l) for l in range(97, 123, 1)]
    upper = [chr(u) for u in range(65, 91, 1)]
    return lower, upper


# -----------------------------------------------------------------------------------------
#  정규식 작업
def reg_raw_text(raws):
    
    raw_results = []
    
    # 웹사이트 주소 지우기 - http:// www.ansancitytour.com -> 띄어쓰기가 되어있으면 안잡힘
    # web_regex = re.compile(r'(http|ftp)s?:\/\/?(www\.)[-a-zA-Z0-9@:%._\+~#=]{2,256}\.[a-z]{2,4}([-a-zA-Z0-9@:%_\+.~#?&//=]*)|(www\.)?(?!ww)[-a-z0-9@:%._\+~#=]{2,256}\.[a-z]{2,4}([-a-zA-Z0-9@:%_\+.~#?&//=])/)?', re.IGNORECASE)

    # 전화번호 지우기 - 021231234(5), (02(2)-123(3)-1234), (02)988-9873, 1234-1234 
    phone_regex = re.compile(r'[\[<{\(]{0,1}(([\[<{\(]{0,1}[\☎\☏]{0,1}[\s]{0,1}[0-9]{0,3}[}>\)\]]{0,1}\-{0,1}[0-9]{3,4}\-{0,1}[0-9]{3,4})|([0-9]{9,10}))[}>\)\]]{0,1}')

    # 시간 지우기 - (10:00, 11:00, 14:00, 15:00, 16:00) / 1:00 ~ 12:00
    time_regex = r'[\[<{\(]{0,1}([0-9]|0[0-9]|1[0-9]|2[0-3]):[0-5][0-9][,\-~\s]{0,2}[}>\)\]]{0,1}'

    # 날짜 지우기 - yyyy.mm.dd | mm.dd.yyyy | dd.mm.yyyy
    date_regex = r'((\d{4})([\s./-]{0,3})(0[1-9]|1[012]|[1-9])([\s./-]{0,3})(0[1-9]|[12][0-9]|3[0-1]|[1-9])([\s./-]{0,3}))'\
                '|((0[1-9]|1[012]|[1-9])([\s./-]{0,3})(0[1-9]|[12][0-9]|3[0-1]|[1-9])([\s./-]{0,3})(\d{4})([\s./-]{0,3}))'\
                '|((0[1-9]|[12][0-9]|3[0-1]|[1-9])([\s./-]{0,3})(0[1-9]|1[012]|[1-9])([\s./-]{0,3})(\d{4})([\s./-]{0,3}))'

    # [숫자 + . + 공백 + 단어] -> [숫자 + . + 단어]: 공백제거
    order_regex = r'(\d)\.\s(\w.*)'

    for raw in raws:
        # too interrupt한 아이들 제거
        # pptx_result_pre = re.sub(r'(\s?■\s?)|(❍)|(▶\s*)|(▸\s*)|(»\s*)|(•\s*)', '', pptx_result_pre)
        raw = re.sub(r"(\s?■\s?)|(‘|※|「|」|`|･|ㅇ|□|❍|▶|▸|»|•|©|—|○|\(\)|-|\*|α|β)\s*", "", raw)
        raw = re.sub(r"\d{1,2}\.\s?", "", raw)
        raw = re.sub(r"\s?\|\s?", "", raw)
        raw = re.sub(r"제\d{1,2}장", "", raw)
        raw = re.sub(r"(C|c)h(\.\s)?apter\s?\d{1,3}", "", raw)
        # space bar 너무 많으면 버리기
        raw = re.sub(r"\s{2,}", "", raw)
        #'.'뒤에 스페이스 주기 (혹시 2개 있을수도 있어서 *로 처리)
        raw = re.sub(r"\.\s{1,2}", ". ", raw)
        # 날짜 제거
        raw = re.sub(r"(\d{2,4}(년|\.|/)\s*)\d{1,2}(월|\.|/)\s*\d{1,2}(일|\s*)", ". ", raw)
        # _ (언더바) 삭제
        raw = re.sub(r"_", ". ", raw)
        # 시간 지우기
        raw = remove_regex(time_regex, raw)
        # 날짜 지우기
        raw = remove_regex(date_regex, raw)
        # 상표 지우기
        raw = re.sub(r"(^㉿{1})|(^©{1})\w+", ". ", raw)
        # 숫자만 제일 앞에 있으면 버리기
        raw = re.sub(r"^\d{1,3}\.?", "", raw)
        # 문장과 문장사이에 . 있을때 띄어쓰기가 없다면 띄어쓰기 해줌
        raw = re.sub(r"([ㄱ-ㅣ가-힣|a-zA-Z|0-9])\.([ㄱ-ㅣ가-힣|a-zA-Z|0-9])", r"\1. \2", raw)
        
        raw_results.append(raw)
        
    return raw_results


# -----------------------------------------------------------------------------------------
# pptx 대문자 예외 상황 처리하기 (잠시 보류)
def pptx_upper_lower_to_a_sentence(pptx_results_pre):
    pptx_results = []
    
    # 대문자 예외처리를 위해서
    lower, upper = lower_or_upper()
    
    # 대, 소문자 구분해서 문장으로 만을어주기
    for idx in range(len(pptx_results_pre)):
        stop_lower_idx = 0
        found_upper_idx = 0
        finally_sentenced = ''
        print(pptx_results_pre[idx])
        if pptx_results_pre[idx][:3] == 'ver':
            stop_lower_idx = idx
            for jdx in range(idx, len(pptx_results_pre)):
                if pptx_results_pre[jdx][0] == 'E':
                    found_upper_idx = jdx
                    finally_sentenced =  pptx_results_pre[found_upper_idx] + pptx_results_pre[stop_lower_idx]

        # del text_list[found_upper_idx]
        if idx == stop_lower_idx:
            pptx_results.append(finally_sentenced)
            # print(finally_sentenced)
        elif idx == found_upper_idx:
            continue 
        elif len(pptx_results_pre[idx]) == 1 and pptx_results_pre[idx][0] in upper:
            continue
        else:
            pptx_results.append(pptx_results_pre[idx])
    
    return pptx_results


# -----------------------------------------------------------------------------------------
#  pptx 작업
def raw_pptx_to_dict(pptx_files_list):
    # ------------------------------------------------------------------------------------
    # pptx dict 선언 => 얘를 결국 리턴
    pptx_contents_dict = {}
    i = 0
    # dict에 잘 들어갔는지 test
    pptx_test_print_1 = ''
    pptx_test_print_2 = ''
    
    
    # ---- pptx 추출 시작 ----
    for pptx_file_list in pptx_files_list:
        print(f'{i}번, pptx 추출 시작했습니다.')
        # a single pptx file을 a list에 넣어줄려고
        pptx_results_pre = []

        # pptx 분석 하기위해 list 만들어주기
        prs = Presentation(pptx_file_list)
        
        # 1개의 pptx 분석시작
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                # 전처리 및 라인별 넣어주기
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text == "":
                        continue
                    text_list = sent_tokenize(paragraph.text)
                    for text in text_list:
                        pptx_results_pre.append(text)
                        
        # i = 잘 동작하는지 확인하기 위해서
        i += 1
        print(f'{i}번, sent_tokenize했습니다.')
        
    #         #대문자 예외처리 (유실되는 것들이 있어서 잠시 보류합니다.)
    #         pptx_results = pptx_upper_lower_to_a_sentence(pptx_results_pre)
    #         print(f'{i}번, 대문자 예외를 처리했습니다.')
        pptx_results = reg_raw_text(pptx_results_pre)
        pptx_results = list(filter(None, pptx_results))
        
        if i == 1:
            pptx_test_print_1 = pptx_file_list[50:-5]
        elif i == 2:
            pptx_test_print_2 = pptx_file_list[50:-5]
        
            
        pptx_contents_dict[pptx_file_list[50:-5]] = pptx_results
    
    # test case 1번째 
    print(pptx_contents_dict.get(pptx_test_print_1))
    # test case 2번째
    print(pptx_contents_dict.get(pptx_test_print_2))
    # 둘 개수 맞는지 확인하기(잘 가져왔는지도)
    print(f'done, pptx_contents_dict : is {len(pptx_contents_dict)}')
    print(f'pptx는 {i}개 가져왔어요')
    
    # pptx_content_dict를 아예 반환
    return pptx_contents_dict
    

# -----------------------------------------------------------------------------------------
#  docx 작업
def raw_docx_to_dict(docx_files_list):
    # ------------------------------------------------------------------------------------
    # docx dict 선언 => 얘를 결국 리턴
    docx_contents_dict = {}
    i = 0
    # dict에 잘 들어갔는지 test
    docx_test_print_1 = ''
    docx_test_print_2 = ''

    
    # ---- docx 추출 시작 ----
    for docx_file_list in docx_files_list:
        print(f'{i}번, docx 추출 시작했습니다.')
        # a single pptx file을 a list에 넣어줄려고
        docx_results_pre = []
        
        # docx 한 개씩 가져오는 중
        with open(docx_file_list, "rb") as docx_file:
            result = mammoth.extract_raw_text(docx_file)
            value = result.value  # The raw text
            # sent_tokenize
            text_list= sent_tokenize(value)
            for text in text_list:
                if text == "":
                    continue
                docx_results_pre.append(text)
        
        # i = 잘 동작하는지 확인하기 위해서
        i += 1
        print(f'{i}번, sent_tokenize했습니다.')

        docx_results = reg_raw_text(docx_results_pre)
        docx_results = list(filter(None, docx_results))

        if i == 1:
            docx_test_print_1 = docx_file_list[50:-5]
        elif i == 2:
            docx_test_print_2 = docx_file_list[50:-5]
        

        docx_contents_dict[docx_file_list[50:-5]] = docx_results

    # test case 1번째 
    print(docx_contents_dict.get(docx_test_print_1))
    # test case 2번째
    print(docx_contents_dict.get(docx_test_print_2))
    # 둘 개수 맞는지 확인하기(잘 가져왔는지도)
    print(f'done, docx_contents_dict : is {len(docx_contents_dict)}')
    print(f'docx는 {i}개 가져왔어요')

    # pptx_content_dict를 아예 반환
    return docx_contents_dict
    
# -----------------------------------------------------------------------------------------
#  파일들 불러온 후, 형식에 맞게 파싱 후 각 파일의 dict에 넣고 한 개의 dict로 update
def all_contents_dict():
    
    # ------------------------------------------------------------------------------------
    # pptx dict 만들기
    pptx_contents_dict = raw_pptx_to_dict(pptx_files_list)
    
    # ------------------------------------------------------------------------------------
    # docx dict 만들기
    docx_contents_dict = raw_docx_to_dict(docx_files_list)
    
    # ------------------------------------------------------------------------------------
    # docx와, pptx merging
    all_contents_dict = {}
    all_contents_dict.update(pptx_contents_dict)
    all_contents_dict.update(docx_contents_dict)
    print(f'total length is : {len(all_contents_dict)}')
    print(list(all_contents_dict.keys()))
    
    # ------------------------------------------------------------------------------------
    # log 출력 (파일명 확인하기 위해서)
    completed_log = open(f'/Users/jihyeoh/Desktop/Farm/4_2019한국표준협회/log_stabdard_37' + '.txt', "w+")
    
    return all_contents_dict


# -----------------------------------------------------------------------------------------
#  language 처리하기
def all_contents_to_each_lan():
    # ------------------------------------------------------------------------------------
    # 한글, 영어 구분해서 각각의 dict에 넣어주기
    
    all_contents_dict = all_contents_dict()
    ko_files = {}
    en_files = {}
    

    i = 0
    j = 1
    for key, value in all_contents_dict.items():
        i += 1
        print(key) 
        language = key[-1]
        key_name = key[:-2]
        print(key_name)
        
        if language == '한':
            ko_files.update({key_name : value})
            print(f'{key[:-2]}, {key_name}, {language}')
            completed_log.write(str(i) + '_[한. DONE READING]' + key_name + ' \n')
        elif language == '영':
            en_files.update({key_name : value})
            print(f'{key[:-2]}, {key_name}, {language}')
            completed_log.write(str(i) + '_[영. DONE READING]' + key_name + ' \n')
        elif language == '병':
            completed_log.write(str(i) + '_[병. DONE READING]' + key_name + ' \n')
        else:
            print('지금 잘 안되고 있어요. 코드 수정필요')
            print(f'{key[:-2]}, {key_name}, {language}')
    
    
    print(list(ko_files.keys()))
    print(list(en_files.keys()))
    

    completed_log.write(f'[DONE READING Total] 한:{str(len(en_files))} 영:{str(len(ko_files))}')
    completed_log.close()
    
    return ko_files, en_files


# -----------------------------------------------------------------------------------------
#  작업할 모든 파일 불러오기
def file_import():
    start = timeit.default_timer()
    # -----------------------------------------------------------------------
    # file 다 져오기

    # file_name_lists = [f for f in glob.glob(f'./4.2019한국표준협회/*.*')]
    pptx_files_list = get_filename_list(file_path, '.pptx')
    docx_files_list = get_filename_list(file_path, '.docx')
    # doc_files_list = get_filename_list(file_path, '.doc')
    # ppt_files_list = get_filename_list(file_path, '.ppt')
    # rtf_files_list = get_filename_list(file_path, '.rtf')
    # etc_files_list = []


    # print(file_path + f'\ndocx: {len(docx_name_lists)}, pptx: {len(pptx_name_lists)}, ppt: {len(ppt_name_lists)}, doc: {len(doc_name_lists)}, rtf: {len(rtf_name_lists)}')
    print(file_path + f'pptx: {len(pptx_files_list)}, docx: {len(docx_files_list)}' )
    
    stop = timeit.default_timer()
    print('Running Time: ', stop - start)
    
    return pptx_files_list, docx_files_list

pptx_files_list, docx_files_list= file_import()


# -----------------------------------------------------------------------------------------
#  excel로 output!!!!!

def dict_to_excel():
    start = timeit.default_timer()
    ko_files, en_files = all_contents_to_each_lan()
    # ------------------------------------------------------------------------------------
    # excel 밖으로 빼내기
    workbook = xlsxwriter.Workbook('/Users/jihyeoh/Desktop/Farm/4_2019한국표준협회/4_xlsx/standard37_.xlsx')
    print('excel 빼내기 파일생성 성공!')
    worksheet = workbook.add_worksheet()
    worksheet.write('A1', '특이사항')
    worksheet.write('B1', 'ko')
    worksheet.write('C1', 'en')
    
    row_idx = 2
    
    i = 0

    # ------------------------------------------------------------------------------------
    # excel 글쓰기

    for ko_key, ko_value in ko_files.items():
        for en_key, en_value in en_files.items():
            if ko_key == en_key:
                i += 1
                print(f'{i}-{ko_key}')
                print(f'{i}-{en_key}')
                worksheet.write('A' + str(row_idx), '>'*10 + ko_key)
                row_idx += 1
                for ko, en in zip_longest(ko_value, en_value, fillvalue=' '):
                    ko_index = get_excel_index('ko', row_idx)
                    en_index = get_excel_index('en', row_idx)
                    worksheet.write(ko_index, ko)
                    worksheet.write(en_index, en)
                    row_idx += 1
          
    workbook.close()
    stop = timeit.default_timer()
    
    print('Running Time: ', stop - start)

dict_to_excel()