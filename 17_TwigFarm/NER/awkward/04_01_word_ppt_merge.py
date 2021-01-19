import mammoth
import glob
import re
import os
import xlsxwriter
import timeit
import sys
from itertools import zip_longest
from pptx import Presentation

# help function
# 한글 / 영어 둘다 아닌 경우
def neitherKoNorEn(text):
    neither = re.compile('((?![ㄱ-ㅣ가-힣|a-zA-Z]).)*')
    return bool(neither.fullmatch(text))


# 한글이 포함되지 않은 경우
def mustEnglish(text):
    never_ko = re.compile('((?![ㄱ-ㅣ가-힣]).)*')
    return bool(never_ko.fullmatch(text))


# 영어가 포함되어있는 경우
def isEnglish(text):
    en = re.compile('.*[a-zA-Z]+')
    # return bool(ko.fullmatch(text))
    return bool(en.match(text))


# 한글이 포함되어있는 경우
def isKorean(text):
    ko = re.compile('.*[ㄱ-ㅣ가-힣]+')
    # return bool(ko.fullmatch(text))
    return bool(ko.match(text))



# mustEnglish - false & isKorean - true
# => 한글만 이거나 한글영어 섞여있을때
# 이경우 다시 isEnglish로 다시 check
# false면 한글만 혹은.. 영어없이 숫자나 특수문자만
# true면 영어+한글 혼합

# mustEnglish - true & isKorean - false
# => 영어만

# --------------------------------------------------------------------------------------
# def sentence_tokenizer(text, source_language_code):
#     tokenizer = stanza.Pipeline(lang=source_language_code, processors='tokenize', verbose=False)
#     doc = tokenizer(text)
#     return [sentence.text for sentence in doc.sentences]


def is_skip(text):
    skip = re.compile('^[0-9『』@%#^&~!★●▪:;()/=+*$,._\\\\-\\s]+')
    return bool(skip.fullmatch(text))


def is_page_num(text):
    is_page = re.compile('.[0-9|pP[\\]{}()<>\\s]+')
    return bool(is_page.fullmatch(text))


def get_excel_index(lang, row_idx):
    if lang == 'ko':
        return f'B{str(row_idx)}'
    else:
        return f'C{str(row_idx)}'


# ppt인 파일 가져오기

from pptx import Presentation           
    
# dictionary로 넣어주기
file_lists_ko = {}
file_lists_en = {}
file_lists_mer = {}   


# dictionary로 넣어주기
file_lists_ko = {}
file_lists_en = {}
file_lists_mer = {}   

# file 다 가져오기
def get_filename_list(path, ext):
    return [f for f in glob.glob(path + f"/*{ext}")]

root_path = '/Users/haley/Desktop/Farm/4_2019한국표준협회'
sub_dir = '/4.2019한국표준협회'
file_path = root_path + sub_dir

def file_to_excel():
    start = timeit.default_timer()
    # -----------------------------------------------------------------------
    # file 다 져오기

    # file_name_lists = [f for f in glob.glob(f'./4.2019한국표준협회/*.*')]
    docx_name_lists = get_filename_list(file_path, '.docx')
    # doc, ppt도 있음
    doc_name_lists = get_filename_list(file_path, '.doc')
    pptx_name_lists = get_filename_list(file_path, '.pptx')
    ppt_name_lists = get_filename_list(file_path, '.ppt')
    rtf_name_lists = get_filename_list(file_path, '.rtf')
    # etc_name_lists = []


    print(file_path + f'\ndocx: {len(docx_name_lists)}, pptx: {len(pptx_name_lists)}, ppt: {len(ppt_name_lists)}, doc: {len(doc_name_lists)}, rtf: {len(rtf_name_lists)}')
    
    
    # ------------------------------------------------------------------------------------
    # ppt dict 만들기
    file_lists_pptx = {}
    i = 0

    # dict에 잘 들어갔는지 test
    pptx_test_print_1 = ''
    pptx_test_print_2 = ''
    
    # ppt 추출 시작
    print('ppt 추출 시작')
    for pptx_name_list in pptx_name_lists:
        i += 1
        # 일단 다 긁어오기 위한 list
        pptx_results_pre = []

        # 최종 잘 나눠진 아이를 넣을 1개 pptx의 list
        pptx_results = []
        
        # pptx 분석 하기위해 넣어주기
        prs = Presentation(pptx_name_list)
        
        # 1개의 pptx 분석시작
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text == '':
                        continue
                    pptx_results_pre.append(paragraph.text)


        for pptx_result_pre in pptx_results_pre:
            if type(pptx_result_pre) == 'list':
                for line in pptx_result_pre:
                    pptx_results.append(line.split())
            else:
                pptx_results.append(pptx_result_pre)
        pptx_results = list(filter(None, pptx_results)) 


        if i == 1:
            pptx_test_print_1 = pptx_name_list[15:-5]
        else:
            print(i, end=' ')

        file_lists_pptx[pptx_name_list[15:-5]] = pptx_results

    print(f'done, file_lists_pptx : is {len(file_lists_pptx)}')
    print(file_lists_pptx.get(pptx_test_print_1))
    print(f'pptx는 {i}개 가져왔습')
    
    
    # ------------------------------------------------------------------------------------
    # docx dict 만들기
    print('docx 추출 시작')
    file_lists_docx = {}
    i = 0
    
    # dict에 잘 들어갔는지 test
    docx_test_print_1 = ''
    docx_test_print_2 = ''
    
    # docx 추출 시작
    for docx_name_list in docx_name_lists:
        i += 1
        docx_results = []
    # print(docx_name_list) # 파일명
    # worksheet.write('A' + str(row_idx), ">" * 20 + docx_name_list)
        with open(docx_name_list, "rb") as docx_file:
            result = mammoth.extract_raw_text(docx_file)
            text = result.value  # The raw text
            contents = text.split('\n |" "')
    
        for line in contents:
            line.replace(".", ". ")
        docx_results = text.split('\n')
        docx_results = list(filter(None, docx_results))

        if i == 1:
            docx_test_print_1 = docx_name_list[15:-5]
        elif i == 2:
            docx_test_print_2 = docx_name_list[15:-5]
        else:
            print(i, end=' ')

        file_lists_docx[docx_name_list[15:-5]] = docx_results

    print(f'done, file_lists_docx : is {len(file_lists_docx)}')
    print(file_lists_docx.get(docx_test_print_1))
    print(f'docx는 {i}개 가져왔습')
    
    
    
    
    # ------------------------------------------------------------------------------------
    # docx와, pptx merging
    file_lists_all = {}
    file_lists_all.update(file_lists_pptx)
    file_lists_all.update(file_lists_docx)
    print(f'total length is : {len(file_lists_all)}')
    # print(list(file_lists_all.keys()))
    
    
    # ------------------------------------------------------------------------------------
    # 한글, 영어 구분해서 각각의 dict에 넣어주기
    ko_files = {}
    en_files = {}
    


    # ------------------------------------------------------------------------------------
    # log 출력 (파일명 확인하기 위해서)
    completed_log = open(f'/Users/haley/Desktop/Farm/4_2019한국표준협회/log_stabdard_3' + '.txt', "w+")
    

    
    i = 0
    j = 1
    for key, value in file_lists_all.items():
        i += 1 
        language = key[-3:]
        key_name = key[35:-4]
        print(key[-3:])
        
        if language == '한':
            ko_files.update({key_name : value})
            print(key_name)
            completed_log.write(str(i) + '_[한. DONE READING]' + key_name + ' \n')
        elif language == '영':
            en_files.update({key_name : value})
            print(key_name)
            completed_log.write(str(i) + '_[영. DONE READING]' + key_name + ' \n')
        elif language == '병':
            j += 1 
        
    
    print(list(ko_files.keys()))
    print(list(en_files.keys()))
    
    completed_log.write(f'[DONE READING Total] 한:{str(len(en_files))} 영:{str(len(ko_files))}')
    completed_log.close()

    
    # ------------------------------------------------------------------------------------
    # excel 밖으로 빼내기
    workbook = xlsxwriter.Workbook('/Users/haley/Desktop/Farm/4_2019한국표준협회/4_xlsx/standard13_.xlsx')
    print('excel 빼내기 파일생성 성공!')
    worksheet = workbook.add_worksheet()
    worksheet.write('B1', 'ko')
    worksheet.write('C1', 'en')
    
    row_idx = 2
    
    i = 0
    

    # ------------------------------------------------------------------------------------
    # excel 글쓰기

    for ko_key, ko_value in ko_files.items():
        for en_key, en_value in en_files.items():
            if ko_key == en_key:
                i += 1
                print(f'{i}-{ko_key}')
                worksheet.write('A' + str(row_idx), '>'*10 + ko_key)
                row_idx += 1
                for ko, en in zip_longest(ko_value, en_value, fillvalue=' '):
                    ko_index = get_excel_index('ko', row_idx)
                    en_index = get_excel_index('en', row_idx)
                    worksheet.write(ko_index, ko)
                    worksheet.write(en_index, en)
                    row_idx += 1
          
    workbook.close()
    stop = timeit.default_timer()
    
    print('Running Time: ', stop - start)

file_to_excel()
