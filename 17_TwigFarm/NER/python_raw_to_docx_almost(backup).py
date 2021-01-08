# -----------------------------------------------------------------------------------------
# library import!!!
import mammoth
import glob
import re
import os
import xlsxwriter
import timeit
import sys
from itertools import zip_longest
from pptx import Presentation
from nltk.tokenize import sent_tokenize
import subprocess
import nltk
# ppt인 파일 가져오기
from pptx import Presentation    




# -----------------------------------------------------------------------------------------
# pattern의 regex를 text에서 찾아 지움
def remove_regex(pattern, text):
    if re.search(pattern, text):
        return re.sub(pattern, '', text)
    return text



# -----------------------------------------------------------------------------------------
# excel numbering
def get_excel_index(lang, row_idx):
    if lang == 'ko':
        return f'B{str(row_idx)}'
    else:
        return f'C{str(row_idx)}'



# -----------------------------------------------------------------------------------------
# file 다 가져오기
def get_filename_list(path, ext):
    print('Change encoding files.......\ \n')
    result = []
    for f in glob.glob(file_path + f"/*{ext}"):
        # run convmv shell script -> file normalization NFC -> NFD (한글자소분리해결)
        subprocess.run(['/usr/local/bin/convmv', '-f', 'utf-8', '-t', 'utf-8', '--nfc', '--notest', f])
        result.append(f)
    return result

root_path = '/Users/jihyeoh/Lexcode/팀 채널 - 인공지능 학습 DB 구축 채널/1. 원본DB(렉스코드)'
sub_dir = '/4.2019한국표준협회'
file_path = root_path + sub_dir



# -----------------------------------------------------------------------------------------
# 대소문자 구별하기 위한 
def lower_or_upper():
    lower = [chr(l) for l in range(97, 123, 1)]
    upper = [chr(u) for u in range(65, 91, 1)]
    return lower, upper



# -----------------------------------------------------------------------------------------
#  작업할 모든 파일 불러오기
def file_import():
    start = timeit.default_timer()
    # -----------------------------------------------------------------------
    # file 다 져오기

    # file_name_lists = [f for f in glob.glob(f'./4.2019한국표준협회/*.*')]
    pptx_files_list = get_filename_list(file_path, '.pptx')
    docx_files_list = get_filename_list(file_path, '.docx')
    # doc_files_list = get_filename_list(file_path, '.doc')
    # ppt_files_list = get_filename_list(file_path, '.ppt')
    # rtf_files_list = get_filename_list(file_path, '.rtf')
    # etc_files_list = []

    print(file_path + f'pptx: {len(pptx_files_list)}, docx: {len(docx_files_list)}' )
    
    return pptx_files_list, docx_files_list

    stop = timeit.default_timer()
    print('Running Time: ', stop - start)

# pptx_files_list, docx_files_list= file_import()

# -----------------------------------------------------------------------------------------
#  pattern의 regex를 text에서 찾아 지움
def remove_regex(pattern, text):
    if re.search(pattern, text):
        return re.sub(pattern, '', text)
    return text



# -----------------------------------------------------------------------------------------
#  정규식 작업
def reg_raw_text(raw):
    
    # 일반 웹사이트 주소 지우기
    web_regex = re.compile(r'((http)s{0,1}://)?[\s]{0,1}[-a-z0-9]{0,3}\.{0,1}[-a-z0-9@:%._\\+~#=]{2,256}\.[a-z]{2,4}', re.IGNORECASE)
    raw = remove_regex(web_regex, raw)

    # 이메일주소 지우기
    email_regex = re.compile(r'([a-z0-9]+[\._]?[a-z0-9]+[@]\w+[.]\w{2,3}.{0,1}\w{0,4})', re.IGNORECASE)
    raw = remove_regex(email_regex, raw)
    
    # 전화번호 지우기 - 021231234(5), (02(2)-123(3)-1234), (02)988-9873, 1234-1234 
    phone_regex = re.compile(r'[\[<{\(]{0,1}(([\[<{\(]{0,1}[\☎\☏]{0,1}[\s]{0,1}[0-9]{0,3}[}>\)\]]{0,1}\-{0,1}[0-9]{3,4}\-{0,1}[0-9]{3,4})|([0-9]{9,10}))[}>\)\]]{0,1}')
    raw = remove_regex(phone_regex, raw)

    # 시간 지우기 - (10:00, 11:00, 14:00, 15:00, 16:00) / 1:00 ~ 12:00
    time_regex = r'[\[<{\(]{0,1}([0-9]|0[0-9]|1[0-9]|2[0-3]):[0-5][0-9][,\-~\s]{0,2}[}>\)\]]{0,1}'
    raw = remove_regex(time_regex, raw)

    # 날짜 지우기 - yyyy.mm.dd | mm.dd.yyyy | dd.mm.yyyy
    date_regex = r'((\d{4})([\s./-]{0,3})(0[1-9]|1[012]|[1-9])([\s./-]{0,3})(0[1-9]|[12][0-9]|3[0-1]|[1-9])([\s./-]{0,3}))'\
                '|((0[1-9]|1[012]|[1-9])([\s./-]{0,3})(0[1-9]|[12][0-9]|3[0-1]|[1-9])([\s./-]{0,3})(\d{4})([\s./-]{0,3}))'\
                '|((0[1-9]|[12][0-9]|3[0-1]|[1-9])([\s./-]{0,3})(0[1-9]|1[012]|[1-9])([\s./-]{0,3})(\d{4})([\s./-]{0,3}))'
    raw = remove_regex(date_regex, raw)
    
    # 영문 들어간 날짜 지우기
    date_regex_2 = r'((January|Feburary|March|April|May|June|July|August|September|October|November|December)|(Jan|Feb|Mar|Apr|Aug|Sep|Oct|Nov|Dec))/s?(\.|,)?\s?/d{1,2}/s?(\.|,)?/d{1,2}'
    raw = re.sub(date_regex_2, '', raw)
    
     # [숫자 + . + 공백 + 단어] -> [숫자 + . + 단어]: 공백제거
    order_regex = r'(\d)\.\s(\w.*)'
    raw = re.sub(order_regex, r"\1.\2", raw)
    
    # 대놓고 특수문자 지워버리기
    raw = re.sub(r"(\s?■\s?)|(‘|※|「|」|`|･|ㅇ|□|❍|▶|▸|»|•|©|—|○|\(\)|-|\*|α|β)\s*", "", raw)
    # index 특수문자 지원버리기 
    raw = re.sub(r'①|②|③|④', '', raw)
    # 알파벳 하나만 덩그러니 있으면 삭제
    raw = re.sub(r"\d{1,2}\.\s?", "", raw)
    # … 제거
    raw = re.sub(r'…+', '', raw)
    # . 숫자 지우기
    raw = re.sub(r'\.\s{1}\d{1,2}', '', raw)
    # . 으로 된 목차 다 지우기
    raw = re.sub(r'\.{2,}\s?\d{1,2}', '', raw)
    # . 다 지우기
    raw = re.sub(r'\.{2,}', '', raw)
    # | 으로 연결된 것들 지우기
    raw = re.sub(r"\s?\|\s?", "", raw)
    # 몇 장 삭제
    raw = re.sub(r"제\d{1,2}장", "", raw)
    # 챕터 삭제
    raw = re.sub(r"(C|c)h(\.\s)?apter\s?\d{1,3}", "", raw)
    # space bar 너무 많으면 버리기
    raw = re.sub(r"\s{2,}", "", raw)
    #'.'뒤에 스페이스 주기 (혹시 2개 있을수도 있어서 *로 처리)
    raw = re.sub(r"\.\s{1,2}", ". ", raw)
    # 날짜 제거
    raw = re.sub(r"(\d{2,4}(년|\.|/)\s*)\d{1,2}(월|\.|/)\s*\d{1,2}(일|\s*)", ". ", raw)
    # 상표 지우기
    raw = re.sub(r"(㉿|®|©){1}\w+", ". ", raw)
    # 숫자만 제일 앞에 있으면 버리기
    raw = re.sub(r"^\d{1,3}\.?", "", raw)
    # 문장과 문장사이에 . 있을때 띄어쓰기가 없다면 띄어쓰기 해줌
    raw = re.sub(r"([ㄱ-ㅣ가-힣|a-zA-Z|0-9])\.([ㄱ-ㅣ가-힣|a-zA-Z|0-9])", r"\1. \2", raw)
    # 전화번호(T+, M+ Mobile)
    raw = re.sub(r'(	\+)?(T|M|Fax|FAX){1}(T \+|Fax	\+|Fax	\+|T	\+)?/d{1,}/s?/d{1,/s?(/d{1,3})?', '', raw)
    # 눈에 띄었던 예외()
    raw = re.sub(r'(level){1}\s{1}\d{1}', '', raw)
    # 숫자가 무더기로있는 경우
    raw = re.sub(r'^\(?\d{1,3}(\s?\d{3})?\)?', '', raw)
    return raw



# -----------------------------------------------------------------------------------------
# pptx 대문자 예외 상황 처리하기 (잠시 보류)
def pptx_upper_lower_to_a_sentence(pptx_results_pre):
    pptx_results = []
    
    # 대문자 예외처리를 위해서
    # lower, upper = lower_or_upper()
    
    # 대, 소문자 구분해서 문장으로 만을어주기
    for idx in range(len(pptx_results_pre)):
        stop_lower_idx = 0
        found_upper_idx = 0
        finally_sentenced = ''
        print(pptx_results_pre[idx])
        if pptx_results_pre[idx][:2] == 'ver':
            stop_lower_idx = idx
            for jdx in range(idx, len(pptx_results_pre)):
                if pptx_results_pre[jdx] == 'E':
                    found_upper_idx = jdx
                    finally_sentenced =  pptx_results_pre[found_upper_idx] + pptx_results_pre[stop_lower_idx]
                    
                     
        # del text_list[found_upper_idx]
        if idx == stop_lower_idx:
            pptx_results.append(finally_sentenced)
            # print(finally_sentenced)
        elif len(pptx_results_pre[idx]) == 1 and pptx_results_pre[idx][0] == 'E':
            continue
        elif idx == found_upper_idx:
            continue
        else:
            pptx_results.append(pptx_results_pre[idx])
    
    return pptx_results



# -----------------------------------------------------------------------------------------
#  pptx 작업
def raw_pptx_to_dict(pptx_files_list):
    # ------------------------------------------------------------------------------------
    # pptx dict 선언 => 얘를 결국 리턴
    pptx_contents_dict = {}
    i = 0
    # dict에 잘 들어갔는지 test
    pptx_test_print_1 = ''
    pptx_test_print_2 = ''
    
    # lower, upper = lower_or_upper()
    # ---- pptx 추출 시작 ----
    for pptx_file_list in pptx_files_list:
        # i = 잘 동작하는지 확인하기 위해서
        i += 1
        
        print(f'{i}번, pptx 추출 시작했습니다.')
        # a single pptx file을 a list에 넣어줄려고
        pptx_results_pre = []

        # pptx 분석 하기위해 list 만들어주기
        prs = Presentation(pptx_file_list)
        
        # 1개의 pptx 분석시작
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                # 전처리 및 라인별 넣어주기
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.text = reg_raw_text(paragraph.text)
                    if paragraph.text == "" or paragraph.text == " " or paragraph.text == ", " or paragraph.text == ","  or paragraph.text == "/" or paragraph.text == ' /' or paragraph.text == "." or paragraph.text == '. ':
                        continue
                    pptx_results_pre.append(paragraph.text)
                    
                # sent_tokenize해주기 위한 전 단계
                pptx_results = []   
                for pptx_result_pre in pptx_results_pre:
                    text_list = sent_tokenize(pptx_result_pre)
                    for text in text_list:
                        if len(text) == 1 and text == 'E':
                            text = text + '→'
                        if text[:2] in 'ver':
                            text = '←'+ text
                        pptx_results.append(text)
                     
        print(f'{i}번, sent_tokenize했습니다.')
        
        # 대문자 예외처리 (유실되는 것들이 있어서 잠시 보류합니다.)
#         pptx_results = pptx_upper_lower_to_a_sentence(pptx_results_pre)
#         print(f'{i}번, 대문자 예외를 처리했습니다.')

        
        pptx_results = list(filter(None, pptx_results))
        
        if i == 1:
            pptx_test_print_1 = pptx_file_list[50:-5]
        elif i == 2:
            pptx_test_print_2 = pptx_file_list[50:-5]
        
            
        pptx_contents_dict[pptx_file_list[50:-5]] = pptx_results
    
    # test case 1번째 
    print(pptx_contents_dict.get(pptx_test_print_1))
    # test case 2번째
    print(pptx_contents_dict.get(pptx_test_print_2))
    # 둘 개수 맞는지 확인하기(잘 가져왔는지도)
    print(f'done, pptx_contents_dict : is {len(pptx_contents_dict)}')
    print(f'pptx는 {i}개 가져왔어요')
    
    # pptx_content_dict를 아예 반환
    return pptx_contents_dict



# -----------------------------------------------------------------------------------------
#  docx 작업
def raw_docx_to_dict(docx_files_list):
    # ------------------------------------------------------------------------------------
    # docx dict 선언 => 얘를 결국 리턴
    docx_contents_dict = {}
    i = 0
    # dict에 잘 들어갔는지 test
    docx_test_print_1 = ''
    docx_test_print_2 = ''

    
    # ---- docx 추출 시작 ----
    for docx_file_list in docx_files_list:
        i += 1
        print(f'{i}번, docx 추출 시작했습니다.')
        # a single pptx file을 a list에 넣어줄려고
        docx_results_pre = []
        
        # docx 한 개씩 가져오는 중
        with open(docx_file_list, "rb") as docx_file:
            result = mammoth.extract_raw_text(docx_file)
            values = result.value  # The raw text
            for value in values:
                if value == "":
                    continue
                docx_results_pre.append(reg_raw_text(value))
        
        # i = 잘 동작하는지 확인하기 위해서
        print(f'{i}번, sent_tokenize했습니다.')
        
        docx_results = []
        for docx_result_pre in docx_results_pre:
            text_list = sent_tokenize(docx_result_pre)
            for text in text_list:
                if text == "":
                    continue
                docx_results.append(text)
        

        docx_results = list(filter(None, docx_results))

        if i == 1:
            docx_test_print_1 = docx_file_list[50:-5]
        elif i == 2:
            docx_test_print_2 = docx_file_list[50:-5]
        

        docx_contents_dict[docx_file_list[50:-5]] = docx_results

    # test case 1번째 
    print(docx_contents_dict.get(docx_test_print_1))
    # test case 2번째
    print(docx_contents_dict.get(docx_test_print_2))
    # 둘 개수 맞는지 확인하기(잘 가져왔는지도)
    print(f'done, docx_contents_dict : is {len(docx_contents_dict)}')
    print(f'docx는 {i}개 가져왔어요')

    # pptx_content_dict를 아예 반환
    return docx_contents_dict
    


# -----------------------------------------------------------------------------------------
#  파일들 불러온 후, 형식에 맞게 파싱 후 각 파일의 dict에 넣고 한 개의 dict로 update
def all_contents_dict():
    
    # ------------------------------------------------------------------------------------
    # pptx dict 만들기
    pptx_contents_dict = raw_pptx_to_dict(pptx_files_list)
    
    # ------------------------------------------------------------------------------------
    # docx dict 만들기
    docx_contents_dict = raw_docx_to_dict(docx_files_list)
    
    # ------------------------------------------------------------------------------------
    # docx와, pptx merging
    all_contents_dict = {}
    all_contents_dict.update(pptx_contents_dict)
    all_contents_dict.update(docx_contents_dict)
    
    print(f'total length is : {len(all_contents_dict)}')
    print(list(all_contents_dict.keys()))
    
    
    # ------------------------------------------------------------------------------------
    # 한글, 영어 구분해서 각각의 dict에 넣어주기
    ko_files = {}
    en_files = {}
    mer_files = {}
    
    # ------------------------------------------------------------------------------------
    # log 출력 (파일명 확인하기 위해서)
    completed_log = open(f'/Users/jihyeoh/Desktop/Farm/4_2019한국표준협회/log_stabdard_49' + '.txt', "w+")
    
    
    # -----------------------------------------------------------------------------------------
    #  language 처리하기
    i = 0
    j = 1
    for key, value in all_contents_dict.items():
        i += 1
        print(key) 
        language = key[-1]
        key_name = key[:-2]
        print(key_name)
        
        if language == '한':
            ko_files.update({key_name : value})
            print(f'{key[:-2]}, {key_name}, {language}')
            completed_log.write(str(i) + '_[한. DONE READING]' + key_name + ' \n')
        elif language == '영':
            en_files.update({key_name : value})
            print(f'{key[:-2]}, {key_name}, {language}')
            completed_log.write(str(i) + '_[영. DONE READING]' + key_name + ' \n')
        elif language == '병':
            completed_log.write(str(i) + '_[병. DONE READING]' + key_name + ' \n')
        else:
            print('지금 잘 안되고 있어요. 코드 수정필요')
            print(f'{key[:-2]}, {key_name}, {language}')
    
    
    print(list(ko_files.keys()))
    print(list(en_files.keys()))
    

    completed_log.write(f'[DONE READING Total] 한:{str(len(en_files))} 영:{str(len(ko_files))}')
    completed_log.close()
    
    return ko_files, en_files



# -----------------------------------------------------------------------------------------
#  excel로 output!!!!!

def dict_to_excel():
    start = timeit.default_timer()
    ko_files, en_files = all_contents_dict()
    # ------------------------------------------------------------------------------------
    # excel 밖으로 빼내기
    workbook = xlsxwriter.Workbook('/Users/jihyeoh/Desktop/Farm/4_2019한국표준협회/4_xlsx/standard49_.xlsx')
    print('excel 빼내기 파일생성 성공!')
    worksheet = workbook.add_worksheet()
    worksheet.write('A1', '특이사항')
    worksheet.write('B1', 'ko')
    worksheet.write('C1', 'en')
    
    row_idx = 2
    
    i = 0

    # ------------------------------------------------------------------------------------
    # excel 글쓰기

    for ko_key, ko_value in ko_files.items():
        ko_key = re.sub(r'(_|-|/s){0,1}', '', ko_key)
        for en_key, en_value in en_files.items():
            en_key = re.sub(r'(_|-|/s){0,1}', '', en_key)
            if ko_key == en_key:
                i += 1
                print(f'{i}-{ko_key}')
                print(f'{i}-{en_key}')
                worksheet.write('A' + str(row_idx), '>'*10 + ko_key)
                row_idx += 1
                for ko, en in zip_longest(ko_value, en_value, fillvalue=' '):
                    ko_index = get_excel_index('ko', row_idx)
                    en_index = get_excel_index('en', row_idx)
                    worksheet.write(ko_index, ko)
                    worksheet.write(en_index, en)
                    row_idx += 1
          
    workbook.close()
    stop = timeit.default_timer()
    
    print('Running Time: ', stop - start)


pptx_files_list, docx_files_list= file_import()
dict_to_excel()