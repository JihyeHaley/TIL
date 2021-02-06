'''
    pptx parser    
    Notes:
    1. 절대경로 말고 상대경로로 맞춰서 사용해주세요.
    2. complete log, error log 추가해주세요.
    3. try, except 사용해주세요.
    4. pptx_to_excel 다시 작성해주세요. parameter는 dict로 받습니다.ß
        pptx_dict = {'한': [filename1, filename2, ..],
                     '영': [filename1, filename2, ..],
                     '병': [filename1, filename2, ..],
                     '양': [filename1, filename2, ..]}
    5. strip_regex_all, pptx_extra_regex 따로 있습니다. utils > regex_functions 확인해주세요.
'''


import re
import xlsxwriter
import timeit
from itertools import zip_longest
from nltk.tokenize import sent_tokenize
from tqdm import tqdm
# ppt인 파일 가져오기
from pptx import Presentation

# 시간찍기
from datetime import datetime

from utils.common_functions import *
from utils.regex_functions import *
from word_pos_extractor import *

def pptx_parser_pre(pptx_file):

    pptx_results_pre = []
    pptx_results = []
    prs = Presentation(pptx_file)

    # ---- pptx 추출 시작 ----
    for slide in prs.slides:
        
        # 1개의 pptx 분석시작
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            # 전처리 및 라인별 넣어주기
            for paragraph in shape.text_frame.paragraphs:
                paragraph.text = regex_cleaner(paragraph.text)
                paragraph.text = pptx_extra_regex(paragraph.text)

                if paragraph.text.strip() in ['.', '/', ',', '', '\'', '\"']:
                    continue
                
                pptx_results_pre.append(paragraph.text)

    # sent_tokenize
    for pptx_result_pre in pptx_results_pre:
        text_list = sent_tokenize(pptx_result_pre)
        for text in text_list:
            pptx_results.append(text)

    pptx_results = list(filter(None, pptx_results))
    return pptx_results



# 파일 이름의 예외로 pair를 못 이룰때를 대비
def check_file_name_exception(file_name):
    file_name = re.sub(r'_', '', file_name)
    file_name = re.sub(r'-', '', file_name)
    file_name = re.sub(r'\s{1}', '', file_name)
    return file_name



def pptx_to_excel(pptx_files_list, sub_path):
    if len(pptx_files_list) == 0:
        print('''-----------------------------------------------------------
        PPTX파일 없습니다.
        ''')
    else:
        print('''-----------------------------------------------------------
        PPTX 작업 시작합니다.
        ''')
        
        start = timeit.default_timer()
        timestamp = datetime.now().strftime("%m%d%H%M")
        ko_i = 0
        error_cnt = 0
        ko_lists = []
        
        which_not_files = pptx_files_list
        which_in_files = []
        
        for pptx_list in pptx_files_list:
            if pptx_list[-6] == '한':
                # print(pptx_list)
                ko_lists.append(pptx_list)

        print(" Total PPTX_KOR: ", len(ko_lists))
        
        workbook = xlsxwriter.Workbook('./results/'  + sub_path  + '/' + sub_path  +  '_pptx_'  + timestamp +'.xlsx') # _mustbessossc
        worksheet = workbook.add_worksheet()
        worksheet.write('A1', 'PATH')
        worksheet.write('B1', 'Raw Data')
        worksheet.write('C1', 'KOR')
        worksheet.write('D1', 'ENG')
        
        row_idx = 2
        total_cnt = 0

        completed_log = open(f'./results/'  + sub_path  + '/' + sub_path + '_log_pptx_' + timestamp + '.txt', "w+")
        pptx_word_list_friend = list()
        
        for idx, ko_list in enumerate(ko_lists):
            try:
                kor_sent_list = pptx_parser_pre(ko_list)
                
                # 중복 제거
                kor_sent_list_full = list()
                for sent in kor_sent_list:
                    if sent not in kor_sent_list_full:
                        kor_sent_list_full.append(sent)

                # 총 몇줄인지 확인
                total_cnt += len(kor_sent_list_full)

                for idx, raw_sent in enumerate(kor_sent_list_full):
                    # 한글, 영어가 같이 있는게 아니라면 건너뛰기
                    if isSentKoreanAndEnglish(raw_sent) == False:
                        continue
                    
                    # A. Path 쓰기
                    a_idx =excel_index_creator('A', row_idx)
                    path = pptx_list.split('/')
                    path =  f'{path[-3]}/{path[-2]}/{path[-1]}/{idx}'
                    worksheet.write(a_idx, path)

                    # raw _sent 형태소 분석 시작
                    # B. Raw Sent 쓰기
                    b_idx =excel_index_creator('B', row_idx)
                    worksheet.write(b_idx, raw_sent)
                    ko_words, en_words= find_pattern_show_words(raw_sent)
                    # print('word_matched: ', word_matched)

                    
                    for j in range(len(ko_words)):
                        # 중복 방지
                        if [ko_words[j], en_words[j]] in pptx_word_list_friend:
                            continue
                        else:
                            pptx_word_list_friend.append([ko_words[j], en_words[j]])
                            # D의 개수가 1개면 skip
                            en_words[j] = en_words[j].strip(' ')
                            if skip_mored_word(en_words[j]) == True:
                                continue

                            else:
                                # C.  ko_word 쓰기
                                c_idx =excel_index_creator('C', row_idx)
                                worksheet.write(c_idx, ko_words[j])


                                # D.  en_word 쓰기
                                d_idx = excel_index_creator('D', row_idx)
                                worksheet.write(d_idx, en_words[j])
                                

                                row_idx += 1

                # Write Complete log
                path = pptx_list.split('/')
                path =  f'{path[-3]}/{path[-2]}/{path[-1]}/{idx}'
                completed_log.write('[DONE READING]' + path + '\n')


            

            except Exception as e:
                print('[ERROR MESSAGE]' + str(e) + '\n')
                completed_log.write('[ERROR MESSAGE]' + pptx_list + str(e) + '\n')
                error_cnt += 1

        # 안써지는거 확인하기
        for idx in which_in_files:
            for jdx in which_not_files:
                if idx == jdx:
                    which_not_files.remove(jdx)
        
        
        workbook.close()
        completed_log.close()
        stop = timeit.default_timer()

        
        print(' PPTX =====> Raw text to excel DONE (Running Time: ', stop - start, "sec.)")
        print(f'total_count = {total_cnt}')