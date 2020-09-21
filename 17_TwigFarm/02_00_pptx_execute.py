import glob

file_name_lists = [f for f in glob.glob(f'./4.2019한국표준협회/*.*x')]
file_name_lists_doc = [f for f in glob.glob(f'./4.2019한국표준협회/*.doc')]
docx_name_lists = []
pptx_name_lists = []

from pptx import Presentation

file_lists_pptx = {}
    
for pptx_name_list in pptx_name_lists:
    # 1개 pptx의 list
    pptx_results = []
    
    # pptx 분석 하기위해 넣어주기
    prs = Presentation(pptx_name_list)
    
    # 1개의 pptx 분석시작
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                pptx_results.append(paragraph.text)
            
        
    file_lists_pptx[pptx_name_list[15:]] = pptx_results

print(f'done, file_lists_pptx : is {len(file_lists_pptx)}')



print(file_lists_pptx)
# # 개수확인
# print(f'Total file: {len(docx_name_lists)}, {len(pptx_name_lists)}')

# file_lists_pptx = {}

# from pptx import Presentation

# # ppt 실행
# for pptx_name_list in pptx_name_lists:
#     # 1개 pptx의 list
#     pptx_results = []
    
#     # pptx 분석 하기위해 넣어주기
#     prs = Presentation(pptx_name_list)
    
#     # 1개의 pptx 분석시작
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if not shape.has_text_frame:
#                 continue
#             for paragraph in shape.text_frame.paragraphs:
#                 pptx_results.append(paragraph.text)
           
    
#     file_lists_pptx[pptx_name_list] = pptx_results

    