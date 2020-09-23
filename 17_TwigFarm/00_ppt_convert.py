import os

path = '/Users/haley/Desktop/Git/TIL/17_TwigFarm/ppt'
file_list = os.listdir(path)


from pptx import Presentation

results1 = []

for file_name_raw in file_list:

    result = []
    result.append(file_name_raw)

    file_name = './ppt/' + file_name_raw
    prs = Presentation(file_name)

    for slide in prs.slides:
        for shape in slide.shapes:
            print(f'type: {type(shape)}, {shape}')
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                print(f'paragraph: {paragraph.text}')
                if paragraph.text == '':
                    continue
                # if  == True:
                #     continue
                result.append(paragraph.text)
    

print(result)