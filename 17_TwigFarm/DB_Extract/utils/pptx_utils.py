import re

from pptx import Presentation


def pptx_extra_regex(raw):
    # 대놓고 특수문자 지워버리기
    raw = re.sub(r'(\s?■\s?)|(‘|※|「|」|`|･|ㅇ|□|❍|▶|▸|»|•|©|—|○|\(\)|-|\*|α|β)\s*', '', raw)
    # index 특수문자 지원버리기
    raw = re.sub(r'①|②|③|④|⑤', '', raw)
    # 알파벳 하나만 덩그러니 있으면 삭제
    raw = re.sub(r"\d{1,2}\.\s?", "", raw)
    # … 제거
    raw = re.sub(r'…+', '', raw)
    # # . 숫자 지우기
    # raw = re.sub(r'\.\s{1}\d{1,2}', '', raw)
    # . 으로 된 목차 다 지우기
    # raw = re.sub(r'\.{2,}\s?\d{1,2}', '', raw)
    # . 다 지우기
    raw = re.sub(r'\.{2,}', '', raw)
    # 주석으로 처리되는 아이들 지우기
    raw = re.sub(r"(ISOfocus|ISO(f포커스|포커스))_?\d{1,3}\s*\|\s*\d{1,2}", "", raw)
    raw = re.sub(r"\s{1}\|\s{1}(ISOfocus|ISO포커스){1}_{1}\d{1,3}", "", raw)
    raw = re.sub(r" – ISO \| ", "", raw)
    # 몇 장 삭제
    raw = re.sub(r"제\d{1,2}장", "", raw)
    # 챕터 삭제
    raw = re.sub(r"(C|c)h(\.\s)?apter\s?\d{1,3}", "", raw)
    # space bar 너무 많으면 버리기
    raw = re.sub(r"\s{2,}", "", raw)
    # '.'뒤에 스페이스 주기 (혹시 2개 있을수도 있어서 *로 처리)
    raw = re.sub(r"\.\s{1,2}", ". ", raw)
    # 날짜 제거
    raw = re.sub(r"(\d{2,4}(년|\.|/)\s*)\d{1,2}(월|\.|/)\s*\d{1,2}(일|\s*)", ". ", raw)
    # 상표 지우기
    raw = re.sub(r"(㉿|®|©){1}\w+", ". ", raw)
    # 숫자만 제일 앞에 있으면 버리기
    raw = re.sub(r"^\d{1,3}\.?", "", raw)

    # 전화번호(T+, M+ Mobile)
    raw = re.sub(r'(	\+)?(T|M|Fax|FAX){1}(T \+|Fax	\+|Fax	\+|T	\+)?/d{1,}/s?/d{1,/s?(/d{1,3})?', '', raw)
    # 눈에 띄었던 예외()
    raw = re.sub(r'(level){1}\s{1}\d{1}', '', raw)
    # 숫자가 무더기로있는 경우
    raw = re.sub(r'^\(?\d{1,3}(\s?\d{3})?\)?', '', raw)

    # 영문 들어간 날짜 지우기
    date_regex_2 = r"((January|Feburary|March|April|May|June|July|August|September|October|November|December)|(Jan|Feb|Mar|Apr|Aug|Sep|Oct|Nov|Dec))/s?(\.|,)?\s?/d{1,2}/s?(\.|,)?/d{1,2}"
    raw = re.sub(date_regex_2, "", raw)

    return raw


def pdf_extra_regex(raw):
    # 대놓고 특수문자 지워버리기
    raw = re.sub(r'(\s?■\s?)|(‘|※|「|」|`|･|ㅇ|□|❍|▶|▸|»|•|©|—|○|\(\)|-|\*|α|β)\s*', '', raw)
    # index 특수문자 지원버리기
    raw = re.sub(r'①|②|③|④|⑤', '', raw)
    # 알파벳 하나만 덩그러니 있으면 삭제
    raw = re.sub(r"\d{1,2}\.\s?", "", raw)
    # … 제거
    raw = re.sub(r'…+', '', raw)
    # # . 숫자 지우기
    # raw = re.sub(r'\.\s{1}\d{1,2}', '', raw)
    # . 으로 된 목차 다 지우기
    # raw = re.sub(r'\.{2,}\s?\d{1,2}', '', raw)
    # . 다 지우기
    raw = re.sub(r'\.{2,}', '', raw)
    # 주석으로 처리되는 아이들 지우기
    raw = re.sub(r"(ISOfocus|ISO(f포커스|포커스))_?\d{1,3}\s*\|\s*\d{1,2}", "", raw)
    raw = re.sub(r"\s{1}\|\s{1}(ISOfocus|ISO포커스){1}_{1}\d{1,3}", "", raw)
    raw = re.sub(r" – ISO \| ", "", raw)
    # 몇 장 삭제
    raw = re.sub(r"제\d{1,2}장", "", raw)
    # 챕터 삭제
    raw = re.sub(r"(C|c)h(\.\s)?apter\s?\d{1,3}", "", raw)
    # space bar 너무 많으면 버리기
    raw = re.sub(r"\s{2,}", "", raw)
    # '.'뒤에 스페이스 주기 (혹시 2개 있을수도 있어서 *로 처리)
    raw = re.sub(r"\.\s{1,2}", ". ", raw)
    # 날짜 제거
    raw = re.sub(r"(\d{2,4}(년|\.|/)\s*)\d{1,2}(월|\.|/)\s*\d{1,2}(일|\s*)", ". ", raw)
    # 상표 지우기
    raw = re.sub(r"(㉿|®|©){1}\w+", ". ", raw)
    # 숫자만 제일 앞에 있으면 버리기
    raw = re.sub(r"^\d{1,3}\.?", "", raw)

    # 전화번호(T+, M+ Mobile)
    raw = re.sub(r'(	\+)?(T|M|Fax|FAX){1}(T \+|Fax	\+|Fax	\+|T	\+)?/d{1,}/s?/d{1,/s?(/d{1,3})?', '', raw)
    # 눈에 띄었던 예외()
    raw = re.sub(r'(level){1}\s{1}\d{1}', '', raw)
    # 숫자가 무더기로있는 경우
    raw = re.sub(r'^\(?\d{1,3}(\s?\d{3})?\)?', '', raw)

    # 영문 들어간 날짜 지우기
    date_regex_2 = r"((January|Feburary|March|April|May|June|July|August|September|October|November|December)|(Jan|Feb|Mar|Apr|Aug|Sep|Oct|Nov|Dec))/s?(\.|,)?\s?/d{1,2}/s?(\.|,)?/d{1,2}"
    raw = re.sub(date_regex_2, "", raw)


def _read_pptx_to_text(pptx_file):

    pptx_results_pre = []
    pptx_results = []
    prs = Presentation(pptx_file)

    # ---- pptx 추출 시작 ----
    for slide in prs.slides:
        
        # 1개의 pptx 분석시작
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            # 전처리 및 라인별 넣어주기
            for paragraph in shape.text_frame.paragraphs:
                paragraph.text = pptx_extra_regex(paragraph.text)

                if paragraph.text.strip() in ['.', '/', ',', '', '\'', '\"']:
                    continue
                
                pptx_results_pre.append(paragraph.text)

    # sent split
    for pptx_result_pre in pptx_results_pre:
        text_list = pptx_result_pre.split('/')
        for text in text_list:
            pptx_results.append(text)

    pptx_results = list(filter(None, pptx_results))
    return pptx_results