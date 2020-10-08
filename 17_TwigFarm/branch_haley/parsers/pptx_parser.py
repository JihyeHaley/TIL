'''
    pptx parser    

    Notes:
    1. 절대경로 말고 상대경로로 맞춰서 사용해주세요.
    2. complete log, error log 추가해주세요.
    3. try, except 사용해주세요.
    4. pptx_to_excel 다시 작성해주세요. parameter는 dict로 받습니다.ß
        pptx_dict = {'한': [filename1, filename2, ..],
                     '영': [filename1, filename2, ..],
                     '병': [filename1, filename2, ..],
                     '양': [filename1, filename2, ..]}
    5. strip_regex_all, pptx_extra_regex 따로 있습니다. utils > regex_functions 확인해주세요.

'''



import xlsxwriter
import timeit
from itertools import zip_longest
from nltk.tokenize import sent_tokenize

# ppt인 파일 가져오기
from pptx import Presentation

from utils.common_functions import *
from utils.regex_functions import *


# 대소문자 구별하기 위한
def lower_or_upper():
    lower = [chr(l) for l in range(97, 123, 1)]
    upper = [chr(u) for u in range(65, 91, 1)]
    return lower, upper


def pptx_to_excel(pptx_lists):

    start = timeit.default_timer()

    try:
        # ---- pptx 추출 시작 ----
        for pptx_file in pptx_lists:
            # i = 잘 동작하는지 확인하기 위해서
            pptx_results_pre = []

            # pptx 분석 하기위해 list 만들어주기
            prs = Presentation(pptx_file)

            # 1개의 pptx 분석시작
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue

                    # 전처리 및 라인별 넣어주기
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.text = pptx_extra_regex(paragraph.text)

                        if paragraph.text.strip() in ['.', '/', ',']:
                            continue
                        pptx_results_pre.append(paragraph.text)

                    # sent_tokenize 해주기 위한 전 단계
                    pptx_results = []
                    for pptx_result_pre in pptx_results_pre:
                        text_list = sent_tokenize(pptx_result_pre)
                        for text in text_list:
                            if len(text) == 1 and text == 'E':
                                text = text + '→'
                            if text[:2] in 'ver':
                                text = '←' + text
                            pptx_results.append(text)

            pptx_results = list(filter(None, pptx_results))


        ko_files, en_files = [], []

        # excel 밖으로 빼내기
        # ------- 절대경로보다 상대경로 사용해주세요
        workbook = xlsxwriter.Workbook('/Users/jihyeoh/Desktop/Farm/4_2019한국표준협회/4_xlsx/standard51_.xlsx')
        worksheet = workbook.add_worksheet()
        worksheet.write('A1', ' ')
        worksheet.write('B1', 'KOR')
        worksheet.write('C1', 'ENG')

        row_idx = 2
        i = 0

        # excel 글쓰기
        for ko_key, ko_value in ko_files.items():
            for en_key, en_value in en_files.items():
                if ko_key == en_key:
                    i += 1
                    print(f'{i}-{ko_key}')
                    print(f'{i}-{en_key}')
                    worksheet.write('A' + str(row_idx), '>' * 10 + ko_key)
                    row_idx += 1
                    for ko, en in zip_longest(ko_value, en_value, fillvalue=' '):
                        worksheet.write('B' + str(row_idx), ko)
                        worksheet.write('C' + str(row_idx), en)
                        row_idx += 1

    except Exception as e:
        print('##')

    workbook.close()
    stop = timeit.default_timer()
    print('Running Time: ', stop - start)

