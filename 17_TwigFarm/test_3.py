import mammoth
import glob
import re
import os
import xlsxwriter
import timeit
import sys
from itertools import zip_longest
from pptx import Presentation
from nltk.tokenize import sent_tokenize
import subprocess
import nltk


# pattern의 regex를 text에서 찾아 지움
def remove_regex(pattern, text):
    if re.search(pattern, text):
        return re.sub(pattern, '', text)
    return text

# 웹사이트 주소 지우기 - http:// www.ansancitytour.com -> 띄어쓰기가 되어있으면 안잡힘
web_regex = re.compile(r'(http|ftp)s?:\/\/?(www\.)[-a-zA-Z0-9@:%._\+~#=]{2,256}\.[a-z]{2,4}([-a-zA-Z0-9@:%_\+.~#?&//=]*)|(www\.)?(?!ww)[-a-z0-9@:%._\+~#=]{2,256}\.[a-z]{2,4}([-a-zA-Z0-9@:%_\+.~#?&//=])', re.IGNORECASE)

# 전화번호 지우기 - 021231234(5), (02(2)-123(3)-1234), (02)988-9873, 1234-1234 
phone_regex = re.compile(r'[\[<{\(]{0,1}(([\[<{\(]{0,1}[\☎\☏]{0,1}[\s]{0,1}[0-9]{0,3}[}>\)\]]{0,1}\-{0,1}[0-9]{3,4}\-{0,1}[0-9]{3,4})|([0-9]{9,10}))[}>\)\]]{0,1}')

# 시간 지우기 - (10:00, 11:00, 14:00, 15:00, 16:00) / 1:00 ~ 12:00
time_regex = r'[\[<{\(]{0,1}([0-9]|0[0-9]|1[0-9]|2[0-3]):[0-5][0-9][,\-~\s]{0,2}[}>\)\]]{0,1}'

# 날짜 지우기 - yyyy.mm.dd | mm.dd.yyyy | dd.mm.yyyy
date_regex = r'((\d{4})([\s./-]{0,3})(0[1-9]|1[012]|[1-9])([\s./-]{0,3})(0[1-9]|[12][0-9]|3[0-1]|[1-9])([\s./-]{0,3}))'\
            '|((0[1-9]|1[012]|[1-9])([\s./-]{0,3})(0[1-9]|[12][0-9]|3[0-1]|[1-9])([\s./-]{0,3})(\d{4})([\s./-]{0,3}))'\
            '|((0[1-9]|[12][0-9]|3[0-1]|[1-9])([\s./-]{0,3})(0[1-9]|1[012]|[1-9])([\s./-]{0,3})(\d{4})([\s./-]{0,3}))'

# [숫자 + . + 공백 + 단어] -> [숫자 + . + 단어]: 공백제거
order_regex = r'(\d)\.\s(\w.*)'




# excel numbering
def get_excel_index(lang, row_idx):
    if lang == 'ko':
        return f'B{str(row_idx)}'
    else:
        return f'C{str(row_idx)}'


# ppt인 파일 가져오기
from pptx import Presentation           
    
# dictionary로 넣어주기
file_lists_ko = {}
file_lists_en = {}
file_lists_mer = {}   


# dictionary로 넣어주기
file_lists_ko = {}
file_lists_en = {}
file_lists_mer = {}   

# file 다 가져오기
def get_filename_list(path, ext):
    print('Change encoding files.......\ \n')
    result = []
    for f in glob.glob(file_path + f"/*{ext}"):
        # run convmv shell script -> file normalization NFC -> NFD (한글자소분리해결)
        subprocess.run(['/usr/local/bin/convmv', '-f', 'utf-8', '-t', 'utf-8', '--nfc', '--notest', f])
        result.append(f)
    return result

root_path = '/Users/jihyeoh/Lexcode/팀 채널 - 인공지능 학습 DB 구축 채널/1. 원본DB(렉스코드)'
sub_dir = '/4.2019한국표준협회'
file_path = root_path + sub_dir

def file_to_excel():
    start = timeit.default_timer()
    # -----------------------------------------------------------------------
    # file 다 져오기

    # file_name_lists = [f for f in glob.glob(f'./4.2019한국표준협회/*.*')]
    docx_name_lists = get_filename_list(file_path, '.docx')
    # doc, ppt도 있음
    doc_name_lists = get_filename_list(file_path, '.doc')
    pptx_name_lists = get_filename_list(file_path, '.pptx')
    ppt_name_lists = get_filename_list(file_path, '.ppt')
    rtf_name_lists = get_filename_list(file_path, '.rtf')
    # etc_name_lists = []


    print(file_path + f'\ndocx: {len(docx_name_lists)}, pptx: {len(pptx_name_lists)}, ppt: {len(ppt_name_lists)}, doc: {len(doc_name_lists)}, rtf: {len(rtf_name_lists)}')
    
    
    # ------------------------------------------------------------------------------------
    # ppt dict 만들기
    file_lists_pptx = {}
    i = 0

    # dict에 잘 들어갔는지 test
    pptx_test_print_1 = ''
    # completed_log = open(f'/Users/jihyeoh/Desktop/Farm/4_2019한국표준협회/log_pptx_3' + '.txt', "w+")

    # ppt 추출 시작
    print('ppt 추출 시작')
    for pptx_name_list in pptx_name_lists:
        i += 1
        print(i, end=' ')
        # 일단 다 긁어오기 위한 list
        pptx_results_pre = []

        
        # pptx 분석 하기위해 넣어주기
        prs = Presentation(pptx_name_list)
        
        # 1개의 pptx 분석시작
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                # 전처리 및 라인별 넣어주기
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.text = sent_tokenize(paragraph.text)
                    if paragraph.text == '' or ', ' or ',' or '.':
                        continue
                    pptx_results_pre.append(paragraph.text)
            
        
        # 최종 잘 나눠진 아이를 넣을 1개 pptx의 list
        pptx_results = []
        
        
        # 전처리 시작
        for pptx_result_pre in pptx_results_pre:
            # too interrupt한 아이들 제거
            # pptx_result_pre = re.sub(r'(\s?■\s?)|(❍)|(▶\s*)|(▸\s*)|(»\s*)|(•\s*)', '', pptx_result_pre)
            pptx_result_pre = re.sub(r'(\s?■\s?)|(ㅇ|□|❍|▶|▸|»|•|©|—|○|\(\)|-|\*|α|β)\s*', '', pptx_result_pre)
            pptx_result_pre = re.sub(r'…*\.*\s?\d{1,3}', '', pptx_result_pre)
            pptx_result_pre = re.sub(r'\d{1,2}\.\s*', '', pptx_result_pre)
            pptx_result_pre = re.sub(r'\s?\|\s?', '', pptx_result_pre)
            # space bar 너무 많으면 버리기
            pptx_result_pre = re.sub(r'\s{2,}', '', pptx_result_pre)
            # /t숫자들 라고 써있는 문자 제거
            pptx_result_pre = re.sub(r'\/t{1}/d{1,}', '', pptx_result_pre)
            #'.'뒤에 스페이스 주기 (혹시 2개 있을수도 있어서 *로 처리)
            pptx_result_pre = re.sub(r'\.\s{1,2}', '. ', pptx_result_pre)
            # 날짜 제거
            pptx_result_pre = re.sub(r'(\d{2,4}(년|\.|/)\s*)\d{1,2}(월|\.|/)\s*\d{1,2}(일|\s*)', '. ', pptx_result_pre)
            # _ (언더바) 삭제
            pptx_result_pre = re.sub(r'_', '. ', pptx_result_pre)
            # URL, EMAIL
            pptx_result_pre = remove_regex(web_regex, pptx_result_pre)
            # 전화번호 지우기
            pptx_result_pre = remove_regex(phone_regex, pptx_result_pre)
            # 시간 지우기
            pptx_result_pre = remove_regex(time_regex, pptx_result_pre)
            # 날짜 지우기
            pptx_result_pre = remove_regex(date_regex, pptx_result_pre)
            #상표 지우기
            pptx_result_pre = re.sub(r'(^㉿{1})|(^©{1})\w+', '. ', pptx_result_pre)
            # 숫자만 제일 앞에 있으면 버리기
            pptx_result_pre = re.sub(r'^\d{1,3}\.?', '', pptx_result_pre)
            # 문장과 문장사이에 . 있을때 띄어쓰기가 없다면 띄어쓰기 해줌
            pptx_result_pre = re.sub(r'([ㄱ-ㅣ가-힣|a-zA-Z|0-9])\.([ㄱ-ㅣ가-힣|a-zA-Z|0-9])', r'\1. \2', pptx_result_pre)                
            
            pptx_results.append(pptx_result_pre)
        pptx_results = list(filter(None, pptx_results)) 
    
        if i == 1:
            pptx_test_print_1 = pptx_name_list[:-5]
        else:
            print(i, end=' ')

        file_lists_pptx[pptx_name_list[:-5]] = pptx_results
        print(file_lists_pptx[pptx_name_list[:-5]])
    #     completed_log.write(str(i) + ' \n' + str(pptx_results) + ' \n ')
    # completed_log.close()
    print(f'done, file_lists_pptx : is {len(file_lists_pptx)}')
    print(file_lists_pptx.get(pptx_test_print_1))
    print(f'pptx는 {i}개 가져왔어요')
    # completed_log.close()
    
    
    # ------------------------------------------------------------------------------------
    # docx dict 만들기
    print('docx 추출 시작')
    file_lists_docx = {}
    i = 0
    
    # dict에 잘 들어갔는지 test
    docx_test_print_1 = ''
    
    # docx 추출 시작
    for docx_name_list in docx_name_lists:
        i += 1
        docx_results = []
    # print(docx_name_list) # 파일명
    # worksheet.write('A' + str(row_idx), ">" * 20 + docx_name_list)
        with open(docx_name_list, "rb") as docx_file:
            result = mammoth.extract_raw_text(docx_file)
            text = result.value  # The raw text
            contents = text.split('\n |" "')
    
        for line in contents:
            line.replace(".", ". ")
        docx_results = text.split('\n')
        docx_results = list(filter(None, docx_results))

        if i == 1:
            docx_test_print_1 = docx_name_list[40:-5]
        else:
            print(i, end=' ')

        file_lists_docx[docx_name_list[40:-5]] = docx_results

    print(f'done, file_lists_docx : is {len(file_lists_docx)}')
    print(file_lists_docx.get(docx_test_print_1))
    print(f'docx는 {i}개 가져왔어요')
    
    
    
    
    # ------------------------------------------------------------------------------------
    # docx와, pptx merging
    file_lists_all = {}
    file_lists_all.update(file_lists_pptx)
    file_lists_all.update(file_lists_docx)
    print(f'total length is : {len(file_lists_all)}')
    # print(list(file_lists_all.keys()))
    
    
    

    # ------------------------------------------------------------------------------------
    # log 출력 (파일명 확인하기 위해서)
    completed_log = open(f'/Users/jihyeoh/Desktop/Farm/4_2019한국표준협회/log_stabdard_27' + '.txt', "w+")
    

    # ------------------------------------------------------------------------------------
    # 한글, 영어 구분해서 각각의 dict에 넣어주기
    ko_files = {}
    en_files = {}
    

    i = 0
    j = 1
    for key, value in file_lists_all.items():
        # i += 1 
        language = key[-1]
        key_name = key[:-5]
        print(key_name)
        
        if language == '한':
            ko_files.update({key_name : value})
            print(f'{key[:-5]}, {key_name}, {language}')
            completed_log.write(str(i) + '_[한. DONE READING]' + key_name + ' \n')
        elif language == '영':
            en_files.update({key_name : value})
            print(f'{key[:-5]}, {key_name}, {language}')
            completed_log.write(str(i) + '_[영. DONE READING]' + key_name + ' \n')
        elif language == '병':
            j += 1 
    
    
    print(list(ko_files.keys()))
    print(list(en_files.keys()))
    

    completed_log.write(f'[DONE READING Total] 한:{str(len(en_files))} 영:{str(len(ko_files))}')
    completed_log.close()

    
    # ------------------------------------------------------------------------------------
    # excel 밖으로 빼내기
    workbook = xlsxwriter.Workbook('/Users/jihyeoh/Desktop/Farm/4_2019한국표준협회/4_xlsx/standard27_.xlsx')
    print('excel 빼내기 파일생성 성공!')
    worksheet = workbook.add_worksheet()
    worksheet.write('B1', 'ko')
    worksheet.write('C1', 'en')
    
    row_idx = 2
    
    i = 0
    

    # ------------------------------------------------------------------------------------
    # excel 글쓰기

    for ko_key, ko_value in ko_files.items():
        for en_key, en_value in en_files.items():
            if ko_key == en_key:
                i += 1
                print(f'{i}-{ko_key}')
                print(f'{i}-{en_key}')
                worksheet.write('A' + str(row_idx), '>'*10 + ko_key)
                row_idx += 1
                for ko, en in zip_longest(ko_value, en_value, fillvalue=' '):
                    ko_index = get_excel_index('ko', row_idx)
                    en_index = get_excel_index('en', row_idx)
                    worksheet.write(ko_index, ko)
                    worksheet.write(en_index, en)
                    row_idx += 1
          
    workbook.close()
    stop = timeit.default_timer()
    
    print('Running Time: ', stop - start)

file_to_excel()
