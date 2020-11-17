'''
    pptx parser    
    Notes:
    1. 절대경로 말고 상대경로로 맞춰서 사용해주세요.
    2. complete log, error log 추가해주세요.
    3. try, except 사용해주세요.
    4. pptx_to_excel 다시 작성해주세요. parameter는 dict로 받습니다.ß
        pptx_dict = {'한': [filename1, filename2, ..],
                     '영': [filename1, filename2, ..],
                     '병': [filename1, filename2, ..],
                     '양': [filename1, filename2, ..]}
    5. strip_regex_all, pptx_extra_regex 따로 있습니다. utils > regex_functions 확인해주세요.
'''


import re
import xlsxwriter
import timeit
from itertools import zip_longest
from nltk.tokenize import sent_tokenize
from tqdm import tqdm
# ppt인 파일 가져오기
from pptx import Presentation

# 시간찍기
from datetime import datetime

from utils.common_functions import *
from utils.regex_functions import *
from word_pos_extractor import *

def pptx_parser_pre(pptx_list):
    # pptx 분석 하기위해 list 만들어주기
    pptx_results_pre = []
    prs = Presentation(pptx_list)
    
    # ---- pptx 추출 시작 ----
    for slide in prs.slides:
        
        # 1개의 pptx 분석시작
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
                
            
            # 전처리 및 라인별 넣어주기
            for paragraph in shape.text_frame.paragraphs:
                # paragraph.text = regex_cleaner(paragraph.text)
                # paragraph.text = pptx_extra_regex(paragraph.text)

                if paragraph.text.strip() in ['.', '/', ',', '', '\'', '\"']:
                    continue
                
                pptx_results_pre.append(paragraph.text)

            # sent_tokenize 
            pptx_results = []
            for pptx_result_pre in pptx_results_pre:
                text_list = sent_tokenize(pptx_result_pre)
                for text in text_list:
                    pptx_results.append(text)
    pptx_results = list(filter(None, pptx_results))
    return pptx_results



# 파일 이름의 예외로 pair를 못 이룰때를 대비
def check_file_name_exception(file_name):
    file_name = re.sub(r'_', '', file_name)
    file_name = re.sub(r'-', '', file_name)
    file_name = re.sub(r'\s{1}', '', file_name)
    return file_name



def pptx_to_excel(pptx_files_list, sub_path):
    if len(pptx_files_list) == 0:
        print('''-----------------------------------------------------------
        PPTX파일 없습니다.
        ''')
    else:
        print('''-----------------------------------------------------------
        PPTX 작업 시작합니다.
        ''')
        
        start = timeit.default_timer()
        timestamp = datetime.now().strftime("%m%d%H%M")
        ko_i = 0
        error_cnt = 0
        ko_lists = []
        
        which_not_files = pptx_files_list
        which_in_files = []
        
        for pptx_list in pptx_files_list:
            if pptx_list[-6] == '한':
                # print(pptx_list)
                ko_lists.append(pptx_list)

        print(" Total PPTX_KOR: ", len(ko_lists))
        
        workbook = xlsxwriter.Workbook('./results/'  + sub_path  + '/' + sub_path  +  '_pptx_'  + timestamp +'.xlsx') # _mustbessossc
        worksheet = workbook.add_worksheet()
        worksheet.write('A1', 'PATH')
        worksheet.write('B1', 'Raw Data')
        worksheet.write('C1', 'KOR')
        worksheet.write('D1', 'ENG')
        # worksheet.write('E1', 'MOR')
        # worksheet.write('F1', '매캡')
        row_idx = 2
        total_cnt = 0

        completed_log = open(f'./results/'  + sub_path  + '/' + sub_path + '_log_pptx_' + timestamp + '.txt', "w+")
        
        for idx, ko_list in enumerate(ko_lists):
            try:
                raw_sents = pptx_parser_pre(ko_list)
                
                # 중복 제거
                raw_sents = list(set(raw_sents))
                # 총 몇줄인지 확인
                total_cnt += len(raw_sents)
                # print(f'{idx} - {len(raw_sents)}')

                for idx, raw_sent in enumerate(raw_sents):
                    # 한글, 영어가 같이 있는게 아니라면 건너뛰기
                    if isSentKoreanAndEnglish(raw_sent) == False:
                        continue
                    
                    # A. Path 쓰기
                    a_idx =excel_index_creator('A', row_idx)
                    path = pptx_list.split('/')
                    path =  f'{path[-3]}/{path[-2]}/{path[-1]}/{idx}'
                    worksheet.write(a_idx, path)

                    # raw _sent 형태소 분석 시작
                    # B. Raw Sent 쓰기
                    b_idx =excel_index_creator('B', row_idx)
                    worksheet.write(b_idx, raw_sent)
                    te, ko_words, en_words, mor_match_list_str = find_pattern_show_words(raw_sent)
                    # print('word_matched: ', word_matched)

                    # F. 형태소 분석되는 세세한 것들 쓰기
                    # f_idx =excel_index_creator('F', row_idx)
                    # worksheet.write(f_idx, te)
                    # print(te)
                    

                    for j in range(len(ko_words)):

                        # D의 개수가 1개면 skip
                        en_words[j] = en_words[j].strip(' ')
                        if skip_mored_word(en_words[j]) == True:
                            continue

                        else:
                            # C.  ko_word 쓰기
                            c_idx =excel_index_creator('C', row_idx)
                            worksheet.write(c_idx, ko_words[j])


                            # D.  en_word 쓰기
                            d_idx = excel_index_creator('D', row_idx)
                            worksheet.write(d_idx, en_words[j])
                            

                            # E.  형태소 패턴 쓰기 (확인을 원할때 사용 print or excel에 작성)
                            # Excel 작성
                            # e_idx = excel_index_creator('E', row_idx)
                            
                            # length가 다를때는 일단 넘어가고 
                            # 형태소 어떤 패턴으로 뽑앗는지 확인하기
                            # print해서 확인
                            # if len(ko_words) != len(mor_match_list_str):
                                # continue
                            # length가 같을때는 쓰게 만들기
                            # worksheet.write(e_idx, mor_match_list_str[j])
                            # print(mor_match_list_str[j])
                            
                            # 다음에 쓰여질 줄을 위해서 row_idx += 1
                            row_idx += 1

                # Write Complete log
                path = pptx_list.split('/')
                path =  f'{path[-3]}/{path[-2]}/{path[-1]}/{idx}'
                completed_log.write('[DONE READING]' + path + '\n')


            

            except Exception as e:
                print('[ERROR MESSAGE]' + str(e) + '\n')
                completed_log.write('[ERROR MESSAGE]' + pptx_list + str(e) + '\n')
                error_cnt += 1

        # 안써지는거 확인하기
        for idx in which_in_files:
            for jdx in which_not_files:
                if idx == jdx:
                    which_not_files.remove(jdx)
        
        
        workbook.close()
        completed_log.close()
        stop = timeit.default_timer()

        
        print(' PPTX =====> Raw text to excel DONE (Running Time: ', stop - start, "sec.)")
        print(f'total_count = {total_cnt}')