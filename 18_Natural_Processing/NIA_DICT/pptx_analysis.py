import glob
import subprocess
from pptx import Presentation    
from mecab import *
import os
import xlsxwriter

lower = [chr(l) for l in range(97, 123, 1)]
upper = [chr(u) for u in range(65, 91, 1)] 

def get_filename_list(path, ext):
    # print('Change encoding files.......\ \n')
    result = []
    for f in glob.glob(file_path + f"/*{ext}"):
        # run convmv shell script -> file normalization NFC -> NFD (한글자소분리해결)
        subprocess.run(['/usr/local/bin/convmv', '-f', 'utf-8', '-t', 'utf-8', '--nfc', '--notest', f])
        result.append(f)
    return result


root_path = '/Users/jihyeoh/Lexcode/팀 채널 - 인공지능 학습 DB 구축 채널/1. 원본DB(렉스코드)'
sub_dir = '/4.2019한국표준협회/'

file_path = root_path + sub_dir
pptx_name_lists = glob.glob('/Users/jihyeoh/Lexcode/팀 채널 - 인공지능 학습 DB 구축 채널/1. 원본DB(렉스코드)/4.2019한국표준협회/TE-한국표준협회-ISOfocus-131-en-한.pptx')
print(pptx_name_lists)

def pptx_to_list(pptx_name_list):
    # if pptx_name_list[-5] == '한':
    prs = Presentation(pptx_name_list)
    pptx_result = []
    for slide in prs.slides:
        for shape in slide.shapes:

            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                print(f'{paragraph.text}')
                if paragraph.text == '':
                    continue
                if paragraph.text.strip() in ['.', '/', ',', '', '\'', '\"']:
                    continue
                pptx_result.append(paragraph.text)
    print(pptx_result)
    return pptx_result

pptx_result = pptx_to_list(pptx_name_lists)
    

def mecab_output(raw_sents):
    workbook = xlsxwriter.Workbook('./의약학_wecab_raw_ko_var_output_test' + '.xlsx') # _mustbessossc
    worksheet = workbook.add_worksheet()
    worksheet.write('A1', 'Raw Sent')
    worksheet.write('B1', 'KOR')
    worksheet.write('C1', 'ENG')
    worksheet.write('D1', 'MOR')
    worksheet.write('E1', '매캡')

    row_idx = 2

    for idx, raw_sent in enumerate(raw_sents):
        # A. Raw Sent 쓰기
        a_idx =excel_index_creator('A', row_idx)
        worksheet.write(a_idx, raw_sent)
        print(f'\n#{idx}---','#'*30,)
        print('A열', raw_sent)

        # raw _sent 형태소 분석 시작
        te, word_matched, mor_match_list_str, mor_match_list = find_pattern_show_words(raw_sent)
        print('word_matched: ', word_matched)
        
        # Raw Sent에 대한 수술 후 뽑혀진 한-영 짝꿍들
        ko_words, en_words = make_str(word_matched)
        print(ko_words, en_words)

        # E. 쓰기
        e_idx =excel_index_creator('E', row_idx)
        worksheet.write(e_idx, te)
        print('E열')
        for j in range(len(ko_words)):
            # B.  ko_word 쓰기
            b_idx =excel_index_creator('B', row_idx)
            worksheet.write(b_idx, ko_words[j])

            # C.  en_word 쓰기
            c_idx = excel_index_creator('C', row_idx)
            worksheet.write(c_idx, en_words[j])


            # D.  en_word 쓰기
            d_idx = excel_index_creator('D', row_idx)
            print(idx, raw_sent, '\n\t', ko_words[j], '-', en_words[j])
            # 한-영 짝꿍이 안 맞으면 엑셀에 아예 raw_sent도 입력이 안되서 
            # length가 다를때는 일단 넘어가고 
            if len(ko_words) != len(mor_match_list_str):
                continue
            # length가 같을때는 쓰게 만들기
            worksheet.write(d_idx, mor_match_list_str[j])

            row_idx += 1

    workbook.close()
## 작업 다 한 후 엑셀파일로 내보내기   
mecab_output(pptx_result)